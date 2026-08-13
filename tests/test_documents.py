import os
import tempfile
import unittest
from pathlib import Path

from document_processor import (
    extract_text,
    is_supported_document,
    _sanitize_filename,
    DocumentError,
)


class SanitizeFilenameTests(unittest.TestCase):
    def test_basic_name(self):
        self.assertEqual(_sanitize_filename("report.pdf"), "report.pdf")

    def test_spaces_replaced(self):
        self.assertEqual(_sanitize_filename("my file.txt"), "my_file.txt")

    def test_special_chars_removed(self):
        self.assertEqual(_sanitize_filename("file@name#1.txt"), "filename1.txt")

    def test_unicode_preserved(self):
        self.assertEqual(_sanitize_filename("документ.pdf"), "документ.pdf")

    def test_empty_falls_back(self):
        self.assertEqual(_sanitize_filename(""), "uploaded_file")

    def test_too_long_truncated(self):
        long_name = "a" * 250 + ".txt"
        result = _sanitize_filename(long_name)
        self.assertLessEqual(len(result), 200)


class SupportedDocumentTests(unittest.TestCase):
    def test_txt_supported(self):
        self.assertTrue(is_supported_document("text/plain", "file.txt"))

    def test_pdf_supported(self):
        self.assertTrue(is_supported_document("application/pdf", "file.pdf"))

    def test_docx_supported(self):
        self.assertTrue(is_supported_document("application/vnd.openxmlformats-officedocument.wordprocessingml.document", "file.docx"))

    def test_xlsx_supported(self):
        self.assertTrue(is_supported_document("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "file.xlsx"))

    def test_pptx_supported(self):
        self.assertTrue(is_supported_document("application/vnd.openxmlformats-officedocument.presentationml.presentation", "file.pptx"))

    def test_csv_supported(self):
        self.assertTrue(is_supported_document("text/csv", "file.csv"))

    def test_json_supported(self):
        self.assertTrue(is_supported_document("application/json", "file.json"))

    def test_xml_supported(self):
        self.assertTrue(is_supported_document("application/xml", "file.xml"))

    def test_html_supported(self):
        self.assertTrue(is_supported_document("text/html", "file.html"))

    def test_zip_supported(self):
        self.assertTrue(is_supported_document("application/zip", "file.zip"))

    def test_unsupported_mime(self):
        self.assertFalse(is_supported_document("application/octet-stream", "file.bin"))

    def test_unsupported_extension(self):
        self.assertFalse(is_supported_document("", "file.xyz"))


class ExtractTextTests(unittest.TestCase):
    def setUp(self):
        self.tmpdir = tempfile.TemporaryDirectory()

    def tearDown(self):
        self.tmpdir.cleanup()

    def test_extract_txt(self):
        p = Path(self.tmpdir.name) / "hello.txt"
        p.write_text("Hello, world!", encoding="utf-8")
        self.assertEqual(extract_text(p), "Hello, world!")

    def test_extract_md(self):
        p = Path(self.tmpdir.name) / "doc.md"
        p.write_text("# Title\n\nParagraph.", encoding="utf-8")
        self.assertIn("Title", extract_text(p))

    def test_extract_json(self):
        p = Path(self.tmpdir.name) / "data.json"
        p.write_text('{"key": "value"}', encoding="utf-8")
        self.assertIn('"key"', extract_text(p))

    def test_extract_csv(self):
        p = Path(self.tmpdir.name) / "data.csv"
        p.write_text("a,b\n1,2\n", encoding="utf-8")
        text = extract_text(p)
        self.assertIn("a,b", text)
        self.assertIn("1,2", text)

    def test_extract_xml(self):
        p = Path(self.tmpdir.name) / "data.xml"
        p.write_text("<root><item>value</item></root>", encoding="utf-8")
        text = extract_text(p)
        self.assertIn("item", text)
        self.assertIn("value", text)

    def test_extract_html(self):
        p = Path(self.tmpdir.name) / "page.html"
        p.write_text("<html><body><h1>Title</h1><p>Text</p></body></html>", encoding="utf-8")
        text = extract_text(p)
        self.assertIn("Title", text)
        self.assertIn("Text", text)

    def test_extract_docx(self):
        try:
            from docx import Document
        except ImportError:
            self.skipTest("python-docx not installed")
        p = Path(self.tmpdir.name) / "doc.docx"
        doc = Document()
        doc.add_paragraph("Hello DOCX")
        doc.save(str(p))
        text = extract_text(p)
        self.assertIn("Hello DOCX", text)

    def test_extract_pdf(self):
        try:
            from pypdf import PdfWriter
        except ImportError:
            self.skipTest("pypdf not installed")
        p = Path(self.tmpdir.name) / "doc.pdf"
        writer = PdfWriter()
        writer.add_blank_page(100, 100)
        with open(p, "wb") as f:
            writer.write(f)
        text = extract_text(p)
        self.assertIn("Page 1", text)

    def test_extract_xlsx(self):
        try:
            from openpyxl import Workbook
        except ImportError:
            self.skipTest("openpyxl not installed")
        p = Path(self.tmpdir.name) / "sheet.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Hello"
        ws["B1"] = "XLSX"
        wb.save(str(p))
        text = extract_text(p)
        self.assertIn("Hello", text)
        self.assertIn("XLSX", text)

    def test_extract_pptx(self):
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")
        p = Path(self.tmpdir.name) / "slides.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Hello PPTX"
        prs.save(str(p))
        text = extract_text(p)
        self.assertIn("Hello PPTX", text)

    def test_extract_zip(self):
        p = Path(self.tmpdir.name) / "archive.zip"
        import zipfile
        with zipfile.ZipFile(p, "w") as zf:
            zf.writestr("file.txt", "hello")
        text = extract_text(p)
        self.assertIn("ZIP archive contents", text)
        self.assertIn("file.txt", text)

    def test_unsupported_raises(self):
        p = Path(self.tmpdir.name) / "file.xyz"
        p.write_text("data", encoding="utf-8")
        with self.assertRaises(DocumentError):
            extract_text(p)


if __name__ == "__main__":
    unittest.main()
