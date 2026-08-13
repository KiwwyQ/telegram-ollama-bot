"""
Document text extraction for the Telegram Ollama bot.

Extracts readable text from uploaded documents so the AI can understand them.
Uses standard library where possible; falls back to common libraries for
complex formats (PDF, DOCX, XLSX, PPTX).

This is a lightweight utility, not a full document-processing framework.
"""
from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Optional
import xml.etree.ElementTree as ET


class DocumentError(Exception):
    """Raised when a document cannot be processed."""


def _read_text_file(path: Path, encoding: str = "utf-8") -> str:
    return path.read_text(encoding=encoding)


def _read_json(path: Path) -> str:
    data = json.loads(path.read_text(encoding="utf-8"))
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_csv(path: Path) -> str:
    out = io.StringIO()
    with path.open("r", encoding="utf-8", newline="") as f:
        reader = csv.reader(f)
        writer = csv.writer(out)
        for i, row in enumerate(reader):
            if i > 500:
                writer.writerow([f"... ({reader.line_num - 1} rows total, truncated)"])
                break
            writer.writerow(row)
    return out.getvalue()


def _read_xml(path: Path) -> str:
    tree = ET.parse(path)
    root = tree.getroot()
    lines = []

    def _dump(elem, indent=0):
        tag = elem.tag
        text = (elem.text or "").strip()
        if text:
            lines.append("  " * indent + f"{tag}: {text}")
        for child in elem:
            _dump(child, indent + 1)

    _dump(root)
    return "\n".join(lines)


class _SimpleHTMLParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        self.parts.append(data)

    def get_text(self) -> str:
        return "".join(self.parts)


def _read_html(path: Path) -> str:
    parser = _SimpleHTMLParser()
    parser.feed(path.read_text(encoding="utf-8"))
    text = parser.get_text()
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise DocumentError("python-docx is not installed. Use '# REQUIRE: python-docx'.") from exc
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        parts.append("[TABLE]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
        parts.append("[/TABLE]")
    return "\n".join(parts)


def _read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentError("pypdf is not installed. Use '# REQUIRE: pypdf'.") from exc
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            parts.append(f"--- Page {i + 1} ---\n{text.strip()}")
    return "\n\n".join(parts)


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise DocumentError("openpyxl is not installed. Use '# REQUIRE: openpyxl'.") from exc
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        rows = list(ws.iter_rows(values_only=True))
        for i, row in enumerate(rows[:200]):
            parts.append(" | ".join(str(c) if c is not None else "" for c in row))
        if len(rows) > 200:
            parts.append(f"... ({len(rows)} rows total, truncated)")
    wb.close()
    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentError("python-pptx is not installed. Use '# REQUIRE: python-pptx'.") from exc
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _read_zip(path: Path) -> str:
    with zipfile.ZipFile(path, "r") as zf:
        names = zf.namelist()
    return "ZIP archive contents:\n" + "\n".join(f"- {name}" for name in names[:100])


# Map MIME types / extensions to extractor functions.
EXTRACTORS = {
    "text/plain": _read_text_file,
    "text/markdown": _read_text_file,
    "application/json": _read_json,
    "text/csv": _read_csv,
    "application/xml": _read_xml,
    "text/xml": _read_xml,
    "text/html": _read_html,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _read_docx,
    "application/pdf": _read_pdf,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _read_xlsx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _read_pptx,
    "application/zip": _read_zip,
    "application/x-zip-compressed": _read_zip,
}

EXTENSION_MAP = {
    ".txt": _read_text_file,
    ".md": _read_text_file,
    ".json": _read_json,
    ".csv": _read_csv,
    ".xml": _read_xml,
    ".html": _read_html,
    ".htm": _read_html,
    ".docx": _read_docx,
    ".pdf": _read_pdf,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
    ".zip": _read_zip,
}


def _sanitize_filename(name: str) -> str:
    name = re.sub(r"[^\w.\- ]+", "", name)
    name = re.sub(r"\s+", "_", name.strip())
    if not name:
        name = "uploaded_file"
    return name[:200]


def extract_text(path: Path, mime_type: str = "") -> str:
    """Extract readable text from a document."""
    ext = Path(path).suffix.lower()
    extractor = EXTRACTORS.get(mime_type) or EXTENSION_MAP.get(ext)
    if not extractor:
        raise DocumentError(f"Unsupported document format: {mime_type or ext}")
    return extractor(path)


def is_supported_document(mime_type: str, filename: str = "") -> bool:
    """Return True if the document type is supported."""
    if mime_type in EXTRACTORS:
        return True
    ext = Path(filename).suffix.lower()
    return ext in EXTENSION_MAP


def get_supported_extensions() -> list[str]:
    return sorted(EXTENSION_MAP.keys())
